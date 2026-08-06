import streamlit as st
import pandas as pd
import io
import zipfile
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication
from PIL import Image

# Imports para geração e manipulação de PDF/PowerPoint
from pptx import Presentation
from reportlab.lib.pagesizes import letter, landscape
from reportlab.pdfgen import canvas
from reportlab.lib import colors

st.set_page_config(page_title="Rokfy — Gerador de Certificados em Massa", layout="wide")

# ==========================================
# 1. ESTILO VISUAL ORIGINAL (PRESERVADO)
# ==========================================
custom_css = """
<style>
    @import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght@300;400;500;600;700;800&display=swap');

    @font-face {
        font-family: 'Chomsky';
        src: local('Chomsky'), local('Chomsky Regular'), local('Chomsky-Regular');
        font-weight: normal; font-style: normal;
    }

    html, body, [class*="css"] {
        font-family: 'Plus Jakarta Sans', -apple-system, sans-serif !important;
        background-color: #FAF6EE !important; color: #1A1A1A !important;
    }
    .stApp { background-color: #FAF6EE !important; }

    [data-testid="stSidebar"] { background-color: #F3ECE0 !important; border-right: 1px solid #E5DBD0 !important; }
    .brand-logo-text { font-family: 'Chomsky', 'UnifrakturMaguntia', serif !important; font-size: 3.8rem !important; color: #E05A47 !important; text-align: center; margin: 0; line-height: 0.9; }
    .brand-tag { font-size: 0.65rem !important; font-weight: 800 !important; letter-spacing: 3px !important; color: #1A1A1A !important; text-transform: uppercase; text-align: center; margin-top: 10px; margin-bottom: 20px; }

    .top-navbar { display: flex; justify-content: space-between; align-items: center; background: #FFFFFF; padding: 14px 28px; border-radius: 14px; box-shadow: 0 4px 15px rgba(0,0,0,0.03); margin-bottom: 25px; }
    .shape-container-top { width: 100%; background-color: #FA8072; border-radius: 24px 24px 0 0; padding: 35px 35px 15px 35px; color: #FFFFFF !important; position: relative; }
    .shape-container-top h1, .shape-container-top p { color: #FFFFFF !important; }
    .wave-divider { width: 100%; height: 60px; margin-bottom: 25px; }
    .wave-divider path { fill: #FA8072; }
    .rokfy-card { background: #FFFFFF; border-radius: 16px; padding: 26px; box-shadow: 0 4px 20px rgba(0,0,0,0.04); border: 1px solid #EAE0D5; margin-bottom: 20px; }
    .stButton>button { background-color: #E05A47 !important; color: #ffffff !important; border: none !important; border-radius: 10px !important; font-weight: 700 !important; padding: 0.65rem 1.6rem !important; }
</style>
"""
st.markdown(custom_css, unsafe_allow_html=True)

def render_header_shape(titulo, sub_titulo):
    st.markdown(
        f"""
        <div class="shape-container-top">
            <h1 style="margin:0;">{titulo}</h1>
            <p style="margin-top:5px; opacity: 0.9;">{sub_titulo}</p>
        </div>
        <svg class="wave-divider" viewBox="0 0 1440 320" preserveAspectRatio="none">
            <path d="M0,96L48,112C96,128,192,160,288,160C384,160,480,128,576,133.3C672,139,768,181,864,186.7C960,192,1056,160,1152,138.7C1248,117,1344,107,1392,101.3L1440,96L1440,0L1392,0C1344,0,1248,0,1152,0C1056,0,960,0,864,0C768,0,672,0,576,0C480,0,384,0,288,0C192,0,96,0,48,0L0,0Z"></path>
        </svg>
        """,
        unsafe_allow_html=True
    )

# ==========================================
# 2. MOTOR DE GERAÇÃO DE CERTIFICADOS
# ==========================================
def gerar_pdf_com_imagem_fundo(imagem_bytes, dados_linha, dados_plataforma, pos_x=400, pos_y=280):
    """Gera um PDF em orientação paisagem usando uma imagem (PNG/JPG) como fundo."""
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=landscape(letter))
    width, height = landscape(letter)

    # Adiciona a imagem de fundo
    if imagem_bytes:
        img_temp = Image.open(io.BytesIO(imagem_bytes))
        img_temp_path = "/tmp/temp_bg.png"
        img_temp.save(img_temp_path)
        c.drawImage(img_temp_path, 0, 0, width=width, height=height)

    # Escreve o texto principal concatenado das colunas da planilha
    c.setFont("Helvetica-Bold", 18)
    c.setFillColor(colors.HexColor("#1A1A1A"))
    
    # Monta texto do certificado
    nome = dados_linha.get("nome", dados_linha.get("Nome", ""))
    texto_principal = f"Certificamos que {nome}"
    c.drawCentredString(pos_x, pos_y, texto_principal)

    # Rodapé / Dados preenchidos na plataforma (CPF, Livro, Registro, etc.)
    c.setFont("Helvetica", 10)
    c.setFillColor(colors.HexColor("#555555"))
    y_extra = 60
    for chave, valor in dados_plataforma.items():
        if valor:
            c.drawString(40, y_extra, f"{chave}: {valor}")
            y_extra -= 14

    c.showPage()
    c.save()
    buffer.seek(0)
    return buffer

def substituir_texto_pptx(pptx_bytes, mapa_substituicao):
    """Substitui tags do tipo {nome}, {cpf} em um modelo do PowerPoint (.pptx)."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for chave, valor in mapa_substituicao.items():
                        tag = f"{{{chave}}}"
                        if tag in paragraph.text:
                            paragraph.text = paragraph.text.replace(tag, str(valor))
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# ==========================================
# 3. INTERFACE E NAVEGAÇÃO
# ==========================================
st.sidebar.markdown('<div class="brand-logo-text">Rokfy</div>', unsafe_allow_html=True)
st.sidebar.markdown('<div class="brand-tag">Certificados em Massa</div>', unsafe_allow_html=True)

menu = st.sidebar.radio("Navegacao", [
    "Gerador em Massa",
    "Instrucoes da Planilha",
    "Configuracao de E-mail (SMTP)"
])

# ------------------------------------------
# TELA 1: GERADOR EM MASSA
# ------------------------------------------
if menu == "Gerador em Massa":
    render_header_shape("Gerador de Certificados em Massa", "Suba o modelo visual, conecte a planilha de participantes e distribua com um clique.")

    st.markdown('<div class="rokfy-card">', unsafe_allow_html=True)
    st.subheader("1. Modelo Visual do Certificado")
    
    tipo_modelo = st.radio("Selecione o formato do seu modelo:", ["Imagem (PNG / JPG / JPEG)", "PowerPoint (.PPTX)"])
    
    arquivo_modelo = st.file_uploader(
        "Upload do Modelo", 
        type=["png", "jpg", "jpeg"] if "Imagem" in tipo_modelo else ["pptx"],
        help="Suba a arte base do seu certificado."
    )
    
    if arquivo_modelo:
        st.success(f"Modelo '{arquivo_modelo.name}' carregado com sucesso!")
    st.markdown('</div>', unsafe_allow_html=True)

    st.markdown('<div class="rokfy-card">', unsafe_allow_html=True)
    st.subheader("2. Upload da Planilha de Participantes")
    
    arquivo_planilha = st.file_uploader(
        "Upload da Planilha (Excel ou CSV)", 
        type=["xlsx", "xls", "csv"],
        help="A planilha DEVE conter obrigatoriamente a coluna 'email' ou 'E-mail'."
    )

    df_dados = None
    coluna_email = None

    if arquivo_planilha:
        try:
            if arquivo_planilha.name.endswith('.csv'):
                df_dados = pd.read_csv(arquivo_planilha)
            else:
                df_dados = pd.read_excel(arquivo_planilha)

            st.write("### Previa da Planilha Carregada:")
            st.dataframe(df_dados.head(5), use_container_width=True)

            # Validação do e-mail obrigatório
            colunas_lowercase = [str(c).strip().lower() for c in df_dados.columns]
            
            if "email" in colunas_lowercase or "e-mail" in colunas_lowercase:
                idx = colunas_lowercase.index("email") if "email" in colunas_lowercase else colunas_lowercase.index("e-mail")
                coluna_email = df_dados.columns[idx]
                st.success(f"Coluna de e-mail obrigatoria identificada: '{coluna_email}'")
            else:
                st.error("A planilha NAO possui uma coluna chamada 'email' ou 'e-mail'. Adicione esta coluna para habilitar o envio.")
        except Exception as e:
            st.error(f"Erro ao ler arquivo de planilha: {e}")
    st.markdown('</div>', unsafe_allow_html=True)

    # Se a planilha e modelo estiverem carregados, exibe dados adicionais e gerador
    if df_dados is not None and coluna_email is not None and arquivo_modelo is not None:
        st.markdown('<div class="rokfy-card">', unsafe_allow_html=True)
        st.subheader("3. Informações Adicionais para Preenchimento na Plataforma")
        st.write("Digite abaixo dados fixos ou complementares que devam constar no documento (ex: Numero do Registro, Livro, Folha, Carga Horaria):")

        col_p1, col_p2, col_p3 = st.columns(3)
        with col_p1:
            num_registro_init = st.text_input("Numero do Registro / Certificado", placeholder="Ex: CERT-2026-001")
            livro_reg = st.text_input("Livro de Registro", placeholder="Ex: Livro 12")
        with col_p2:
            num_ordem = st.text_input("Numero de Ordem", placeholder="Ex: Fls. 45")
            folha_reg = st.text_input("Folha de Registro", placeholder="Ex: 088")
        with col_p3:
            cpf_generico = st.text_input("CPF (se nao estiver na planilha)", placeholder="Opcional")
            txt_extra = st.text_input("Observacoes Adicionais", placeholder="Ex: Portaria do MEC nº 123")

        dados_plataforma = {
            "Registro": num_registro_init,
            "Livro": livro_reg,
            "Ordem": num_ordem,
            "Folha": folha_reg,
            "CPF": cpf_generico,
            "Observacoes": txt_extra
        }

        st.markdown("---")
        st.subheader("4. Processamento e Distribuicao")
        
        tab_download, tab_email = st.tabs(["Baixar Todos em ZIP", "Enviar por E-mail em Massa"])

        # OPÇÃO A: DOWNLOAD EM LOTE (ZIP)
        with tab_download:
            st.write("Gere todos os certificados em PDF compactados em um unico arquivo ZIP.")
            if st.button("Gerar Arquivo ZIP com Todos os Certificados"):
                zip_buffer = io.BytesIO()
                with zipfile.ZipFile(zip_buffer, "a", zipfile.ZIP_DEFLATED, False) as zip_file:
                    for index, row in df_dados.iterrows():
                        nome_dest = str(row.get("nome", row.get("Nome", f"participante_{index+1}"))).strip()
                        dados_linha = row.to_dict()

                        if "Imagem" in tipo_modelo:
                            pdf_bytes = gerar_pdf_com_imagem_fundo(
                                arquivo_modelo.getvalue(),
                                dados_linha,
                                dados_plataforma
                            )
                            zip_file.writestr(f"Certificado_{nome_dest}.pdf", pdf_bytes.getvalue())
                        else:
                            # Caso seja PowerPoint
                            mapa = {**dados_linha, **dados_plataforma}
                            pptx_mod = substituir_texto_pptx(arquivo_modelo.getvalue(), mapa)
                            zip_file.writestr(f"Certificado_{nome_dest}.pptx", pptx_mod.getvalue())

                zip_buffer.seek(0)
                st.success("Todos os certificados foram gerados com sucesso!")
                st.download_button(
                    label="Baixar Arquivo ZIP",
                    data=zip_buffer,
                    file_name="Certificados_Rokfy_Massa.zip",
                    mime="application/zip"
                )

        # OPÇÃO B: ENVIO POR E-MAIL
        with tab_email:
            st.write("Envie automaticamente o certificado em anexo para cada e-mail cadastrado na planilha.")
            
            assunto_email = st.text_input("Assunto do E-mail", value="Seu Certificado de Participacao esta disponivel!")
            corpo_email = st.text_area("Corpo da Mensagem", value="Ola! Segue em anexo o seu certificado oficial. Parabéns!")

            if st.button("Disparar E-mails em Massa"):
                if "smtp_config" not in st.session_state:
                    st.error("Configure as credenciais do Servidor SMTP na aba 'Configuração de E-mail (SMTP)' na barra lateral antes de enviar.")
                else:
                    smtp_data = st.session_state["smtp_config"]
                    sucessos = 0
                    erros = 0

                    progresso = st.progress(0)
                    total = len(df_dados)

                    for index, row in df_dados.iterrows():
                        email_dest = str(row[coluna_email]).strip()
                        nome_dest = str(row.get("nome", row.get("Nome", "Participante"))).strip()
                        
                        try:
                            # Gerar PDF
                            dados_linha = row.to_dict()
                            pdf_file = gerar_pdf_com_imagem_fundo(
                                arquivo_modelo.getvalue(),
                                dados_linha,
                                dados_plataforma
                            )

                            # Montar e-mail
                            msg = MIMEMultipart()
                            msg['From'] = smtp_data['usuario']
                            msg['To'] = email_dest
                            msg['Subject'] = assunto_email
                            msg.attach(MIMEText(corpo_email, 'plain'))

                            # Anexo PDF
                            anexo = MIMEApplication(pdf_file.getvalue(), _subtype="pdf")
                            anexo.add_header('Content-Disposition', 'attachment', filename=f"Certificado_{nome_dest}.pdf")
                            msg.attach(anexo)

                            # Envio SMTP
                            server = smtplib.SMTP(smtp_data['server'], int(smtp_data['porta']))
                            server.starttls()
                            server.login(smtp_data['usuario'], smtp_data['senha'])
                            server.send_message(msg)
                            server.quit()

                            sucessos += 1
                        except Exception as e:
                            erros += 1

                        progresso.progress((index + 1) / total)

                    st.success(f"Disparo concluido! Enviados com sucesso: {sucessos} | Falhas: {erros}")

        st.markdown('</div>', unsafe_allow_html=True)

# ------------------------------------------
# TELA 2: INSTRUÇÕES DA PLANILHA
# ------------------------------------------
elif menu == "Instrucoes da Planilha":
    render_header_shape("Estrutura da Planilha", "Saiba como formatar seu arquivo Excel ou CSV.")
    
    st.markdown('<div class="rokfy-card">', unsafe_allow_html=True)
    st.subheader("Regras da Planilha:")
    st.markdown("""
    1. **E-mail (Obrigatório):**
       * Deve haver uma coluna chamada obrigatoriamente **`email`** ou **`E-mail`**. Ela servira para fazer o disparo dos documentos.
    
    2. **Outras Informações (100% Livres):**
       * Todas as outras colunas sao opcionais e flexiveis. Voce pode adicionar:
         * `Nome` ou `nome`
         * `CPF`
         * `Curso`
         * `Data`
         * `Carga Horaria`
         * `Registro`, `Livro`, `Folha` (se desejar colocar na planilha ao inves da plataforma)

    3. **Substituição de Tags no PowerPoint (.PPTX):**
       * Caso utilize um modelo do PowerPoint, insira textos no seu slide no formato `{nome_da_coluna}`.
       * *Exemplo:* `Certificamos que {Nome}, portador do CPF {CPF}, concluiu o curso de {Curso}.`
    """)
    st.markdown('</div>', unsafe_allow_html=True)

# ------------------------------------------
# TELA 3: CONFIGURAÇÃO DE E-MAIL (SMTP)
# ------------------------------------------
elif menu == "Configuracao de E-mail (SMTP)":
    render_header_shape("Configuracao de E-mail (SMTP)", "Conecte sua conta de e-mail para habilitar disparos em massa.")
    
    st.markdown('<div class="rokfy-card">', unsafe_allow_html=True)
    st.subheader("Dados do Servidor SMTP")
    
    servidor_smtp = st.text_input("Servidor SMTP", value="smtp.gmail.com")
    porta_smtp = st.number_input("Porta SMTP", value=587)
    usuario_smtp = st.text_input("Seu E-mail de Envio")
    senha_smtp = st.text_input("Senha / Senha de Aplicativo", type="password")

    if st.button("Salvar Configuracoes SMTP"):
        if servidor_smtp and usuario_smtp and senha_smtp:
            st.session_state["smtp_config"] = {
                "server": servidor_smtp,
                "porta": porta_smtp,
                "usuario": usuario_smtp,
                "senha": senha_smtp
            }
            st.success("Configuracoes de e-mail salvas na sessao com sucesso!")
        else:
            st.error("Preencha todos os campos do servidor SMTP.")
    st.markdown('</div>', unsafe_allow_html=True)
